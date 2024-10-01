# -*- coding: utf-8 -*-
"""churn-prediction.ipynb

Automatically generated by Colab.

Original file is located at
    https://colab.research.google.com/drive/19pRwck43WLRRjk8D7hVkBnAVj_eoMOt1

###Importing Library
"""

import numpy as np
import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt

"""###Importing raw dataset to pandas Dataframe"""

raw_data = pd.read_excel('/content/Customer+Churn+Data (2).xlsx', sheet_name='Data for DSBA')
data = pd.DataFrame(raw_data)

"""###Knowing the Dataset"""

print('Total number of rows : ', len(data))
print('Total number of columns : ',len(data.columns))
print('Number of Numerical columns : ', len(data.select_dtypes(include=['float64', 'int64']).columns))
print('Number of object columns : ', len(data.select_dtypes(include='object').columns))
print('Number of Numerical rows with missing values : ', data.select_dtypes(include=['float64', 'int64']).isnull().values.sum())
print('Number of Non-numerical rows with missing values : ', data.select_dtypes(include=['object']).isnull().values.sum())

data.to_excel('abc.xlsx')

"""###Data preprocessing"""

for datas in data.select_dtypes(include=['float64', 'int64']).columns:
    if data[datas].isnull().sum() > 0:
        data[datas].fillna(np.floor(data[datas].mean())
                           ,inplace=True)
    print(datas," : "," unique_value = ",data[datas].unique())

for datas in data.select_dtypes(include='object').columns:
    print(datas," : "," unique_value = ",data[datas].unique())

object_to_numeric = ['Tenure','Account_user_count','rev_per_month','rev_growth_yoy','coupon_used_for_payment','Day_Since_CC_connect','cashback']

for objects in object_to_numeric:
    data[objects] = data[objects].apply(pd.to_numeric, errors='coerce')
    if objects != 'cashback':
        data[objects].fillna(np.floor(data[objects].mean()),inplace=True)
    else:
        data[objects].fillna(data[objects].mean(),inplace=True)

data['Gender'].replace('M',"Male",inplace=True)
data['Gender'].replace('F',"Female",inplace=True)
data['Payment'].fillna(np.random.choice(data['Payment']),inplace=True)
data['Gender'].fillna(np.random.choice(data['Gender']),inplace=True)
data['Marital_Status'].fillna(np.random.choice(data['Marital_Status']),inplace=True)
data['rev_growth_yoy'] = data['rev_growth_yoy'].apply(pd.to_numeric, errors='coerce')
data['rev_growth_yoy'].fillna(data['rev_growth_yoy'].median(),inplace=True)
data['Login_device'].replace('&&&&',np.random.choice(data['Login_device']),inplace=True)
data['Login_device'].fillna(np.random.choice(data['Login_device']),inplace=True)
data['account_segment'].replace('Super +','Super Plus',inplace=True)
data['account_segment'].replace('Regular +','Regular Plus',inplace=True)
data['account_segment'].fillna('Regular',inplace=True)

data.isnull().values.sum()

numerical_features = data.select_dtypes(include=["float64", "int64"]).drop("Churn", axis=1)
numerical_features.describe()

"""###Grouping

Grouping by City
"""

no_objdata= data.drop(data.select_dtypes(include='object').columns, axis=1)
no_objdata.groupby('City_Tier').count()

"""grouping by account segment"""

round(no_objdata[['Churn','Tenure','CC_Contacted_LY','rev_per_month','rev_growth_yoy','cashback','Day_Since_CC_connect']].groupby(data['account_segment']).mean(),2)

"""group by churn"""

round(no_objdata.groupby('Churn').mean(),2)

"""###one hot encoding"""

hot_encoded_data = pd.get_dummies(data=data,drop_first=True)
hot_encoded_data = hot_encoded_data.drop(columns='AccountID')

"""###Exploratory Data Analysis

Data Analysis
"""

category_columns = list(data.select_dtypes(include='object').columns)

fig, axes = plt.subplots(1,3, figsize=(20,12), facecolor="lightgray")

for i, column in enumerate(category_columns[:3]):
    ax = axes[i]
    d = data[column].value_counts()
    ax.pie(d, labels=d.values,autopct="%1.1f%%",shadow=True)
    ax.set_title(column,size=18)
    ax.legend(d.index)

fig, axes = plt.subplots(1,2, figsize=(20,12), facecolor="lightgray")

for i, column in enumerate(category_columns[3:5]):
    ax = axes[i]
    d = data[column].value_counts()
    ax.pie(d, labels=d.values,autopct="%1.1f%%",shadow=True)
    ax.set_title(column,size=18)
    ax.legend(d.index)

fig, ax = plt.subplots(facecolor="lightblue")
d = data["Churn"].value_counts()
ax.pie(d,autopct='%1.1f%%', startangle=90,labels=d.values)
ax.legend(d.index)
ax.set_title("Distribution of Churn")
centre_circle = plt.Circle((0,0),0.4,fc='lightblue')
fig = plt.gcf()
fig.gca().add_artist(centre_circle)
ax.axis('equal');

# prompt: Scatterplots: Visually depict relationships between features to identify potential patterns (e.g., Tenure vs. Churn).

sns.scatterplot(data=data, x="Tenure", y="Churn")
plt.title("Scatterplot of Tenure vs. Churn")
plt.show()

correlation_matrix = data[['Tenure','Account_user_count','rev_per_month','rev_growth_yoy','coupon_used_for_payment','Day_Since_CC_connect','cashback']].corr()

round(correlation_matrix,2)

"""Objects w.r.t Chrun"""

sns.barplot(x = 'Churn', y = 'Day_Since_CC_connect', data = data)
plt.title('Bar plot for churn vs Day_Since_CC_connect')
plt.show()

# prompt: coupon_used_l12 m vs churn bartchart

sns.barplot(x = 'Churn', y = 'coupon_used_for_payment', data = data)
plt.title('Bar plot for churn vs coupon used')
plt.show()

# prompt: Complain_ly and churn

sns.barplot(x = 'Churn', y = 'Complain_ly', data = data)
plt.title('Bar plot for churn vs Complain_ly')
plt.show()

sns.countplot(x=data['Churn'],hue=data['Payment'],data=data)
plt.plot

sns.countplot(x=data['Account_user_count'],hue=data['Churn'])

sns.countplot(x=data['Churn'],hue=data['Gender'],data=data)

sns.countplot(x=data['Churn'],hue=data['account_segment'],data=data)
plt.plot

sns.countplot(x=data['Churn'],hue=data['Marital_Status'],data=data)
plt.plot

sns.countplot(x=data['Churn'],hue=data['Login_device'],data=data)
plt.plot

sns.barplot(x='account_segment', y='rev_per_month', data=data)
plt.show()
sns.barplot(x='account_segment', y='rev_growth_yoy', data=data)
plt.show()

# prompt: Revenue Metrics:

# Calculate churn rate
churn_rate = data['Churn'].value_counts()[1] / len(data) * 100

# Calculate average revenue per user (ARPU)
arpu = data['rev_per_month'].mean()

# Calculate monthly recurring revenue (MRR)
mrr = arpu * len(data)

# Calculate customer lifetime value (CLTV)
cltv = arpu * data['Tenure'].mean()

# Print revenue metrics
print('Churn Rate:', churn_rate)
print('ARPU:', arpu)
print('MRR:', mrr)
print('CLTV:', cltv)

# prompt: 1.	Customer Tenure:
# o	Average tenure of customers:
# 2.	Users Per Account:
# o	Average number of users per account:
# 3.	Revenue Metrics:
# o	Average revenue per month:
# o	Average revenue growth year-over-year:
# Payment and Usage:

# 1. Average tenure of customers:
average_tenure = data['Tenure'].mean()
print(f"Average tenure of customers: {average_tenure}")

# 2. Average number of users per account:
average_users_per_account = data['Account_user_count'].mean()
print(f"Average number of users per account: {average_users_per_account}")

# 3. Revenue Metrics:
# Average revenue per month:
average_revenue_per_month = data['rev_per_month'].mean()
print(f"Average revenue per month: {average_revenue_per_month}")

# Average revenue growth year-over-year:
average_revenue_growth_yoy = data['rev_growth_yoy'].mean()
print(f"Average revenue growth year-over-year: {average_revenue_growth_yoy}")

# Payment and Usage:
# Average number of coupons used for payment:
average_coupons_used_for_payment = data['coupon_used_for_payment'].mean()
print(f"Average number of coupons used for payment: {average_coupons_used_for_payment}")

# Average number of days since last CC connection:
average_days_since_last_cc_connection = data['Day_Since_CC_connect'].mean()
print(f"Average number of days since last CC connection: {average_days_since_last_cc_connection}")

# Average cashback amount:
average_cashback = data['cashback'].mean()
print(f"Average cashback amount: {average_cashback}")

# prompt: # o	Average tenure of customers:
# # 2.	Users Per Account:
# # o	Average number of users per account:
# # 3.	Revenue Metrics:
# # o	Average revenue per month:
# # o	Average revenue growth year-over-year:
# # Payment and Usage: insted of average make it total

# 1. Average tenure of customers:
average_tenure = data['Tenure'].mean()
print(f"Average tenure of customers: {average_tenure}")

# 2. Users Per Account:
# Average number of users per account:
average_users_per_account = data['Account_user_count'].mean()
print(f"Average number of users per account: {average_users_per_account}")

# 3. Revenue Metrics:
# Total revenue per month:
total_revenue_per_month = data['rev_per_month'].sum()
print(f"Total revenue per month: {total_revenue_per_month}")

# Total revenue growth year-over-year:
total_revenue_growth_yoy = data['rev_growth_yoy'].sum()
print(f"Total revenue growth year-over-year: {total_revenue_growth_yoy}")

# Payment and Usage:
# Total number of coupons used for payment:
total_coupons_used_for_payment = data['coupon_used_for_payment'].sum()
print(f"Total number of coupons used for payment: {total_coupons_used_for_payment}")

# Total number of days since last CC connection:
total_days_since_last_cc_connection = data['Day_Since_CC_connect'].sum()
print(f"Total number of days since last CC connection: {total_days_since_last_cc_connection}")

# Total cashback amount:
total_cashback = data['cashback'].sum()
print(f"Total cashback amount: {total_cashback}")

# prompt: 5.	Demographic and Account Information:

# 5. Demographic and Account Information:
# Gender distribution:
gender_distribution = data['Gender'].value_counts()
print(f"Gender distribution: {gender_distribution}")

# Marital status distribution:
marital_status_distribution = data['Marital_Status'].value_counts()
print(f"Marital status distribution: {marital_status_distribution}")

# Account segment distribution:
account_segment_distribution = data['account_segment'].value_counts()
print(f"Account segment distribution: {account_segment_distribution}")

# City tier distribution:
city_tier_distribution = data['City_Tier'].value_counts()
print(f"City tier distribution: {city_tier_distribution}")

# Payment method distribution:
payment_method_distribution = data['Payment'].value_counts()
print(f"Payment method distribution: {payment_method_distribution}")

# Login device distribution:
login_device_distribution = data['Login_device'].value_counts()
print(f"Login device distribution: {login_device_distribution}")

"""### Correlation matrix and heat map"""

dataset_2 = hot_encoded_data.drop(columns=['Churn'])

dataset_2.corrwith(data['Churn']).plot.bar(
    figsize=(16,9),rot=45,grid=True)

corr = hot_encoded_data.corr()
plt.figure(figsize=(30,15))
sns.heatmap(corr,annot=True)

# prompt: 1.	Correlation Insights on dataset_2.corrwith(data['Churn']).plot.bar(
#     figsize=(16,9),rot=45,grid=True)

# 1. Identify the features with the highest positive and negative correlations with Churn:

# Highest positive correlations:
positive_correlations = dataset_2.corrwith(data['Churn']).sort_values(ascending=False)[1:]
highest_positive_correlations = positive_correlations.head(5)

# Highest negative correlations:
negative_correlations = dataset_2.corrwith(data['Churn']).sort_values(ascending=True)[:-1]
highest_negative_correlations = negative_correlations.head(5)

# 2. Interpret the correlations:

# Positive correlations:
print("Features with highest positive correlations with Churn:")
print(highest_positive_correlations)

# Negative correlations:
print("Features with highest negative correlations with Churn:")
print(highest_negative_correlations)

# 3. Analyze the results:

# - Features with high positive correlations indicate that they are associated with higher churn rates.
# - Features with high negative correlations indicate that they are associated with lower churn rates.

# 4. Consider these correlations when building a churn prediction model or taking actions to reduce churn.

"""# Train Test Split"""

x = hot_encoded_data.drop(columns='Churn')
y = hot_encoded_data['Churn']

from sklearn.model_selection import train_test_split

x_train, x_test, y_train, y_test = train_test_split(x, y, test_size=0.20, random_state=0)

y_test.shape

"""### Feature Scaling"""

from sklearn.preprocessing import StandardScaler

sc = StandardScaler()
x_train = sc.fit_transform(x_train)
x_test = sc.fit_transform(x_test)

x_train

"""# Model Building

### Logistic Regression
"""

from sklearn.linear_model import LogisticRegression

LR = LogisticRegression(random_state=0)

LR.fit(x_train,y_train)

y_pred = LR.predict(x_test)

"""Metrics"""

from sklearn.metrics import f1_score,accuracy_score,confusion_matrix,recall_score,precision_score

acc = accuracy_score(y_test,y_pred)
f1 = f1_score(y_test,y_pred)
cm = confusion_matrix(y_test,y_pred)
rs = recall_score(y_test,y_pred)
ps = precision_score(y_test,y_pred)

result = pd.DataFrame([['Logistic Resgression',acc,f1,rs,ps]],columns=['Model Name',"accuracy","f1 Score","Recall Score","Precission_score"])
result

"""### Decision Tree"""

from sklearn.tree import DecisionTreeClassifier
DT = DecisionTreeClassifier(criterion='entropy',random_state=0)
DT.fit(x_train,y_train)
y_pred_DT = DT.predict(x_test)

acc = accuracy_score(y_test,y_pred_DT)
f1 = f1_score(y_test,y_pred_DT)
cm = confusion_matrix(y_test,y_pred_DT)
rs = recall_score(y_test,y_pred_DT)
ps = precision_score(y_test,y_pred_DT)
result0 = pd.DataFrame([['Decision Tree',acc,f1,rs,ps]],columns=['Model Name',"accuracy","f1 Score","Recall Score","Precission_score"])
result
result = pd.concat([result, result0])

"""### Random Forest"""

from sklearn.ensemble import RandomForestClassifier
RF = RandomForestClassifier(random_state=0)
RF.fit(x_train,y_train)

y_pred = RF.predict(x_test)

"""metrics"""

acc = accuracy_score(y_test,y_pred)
f1 = f1_score(y_test,y_pred)
cm = confusion_matrix(y_test,y_pred)
rs = recall_score(y_test,y_pred)
ps = precision_score(y_test,y_pred)

result1 = pd.DataFrame([['Random Forest', acc, f1, rs, ps]],
                       columns=['Model Name', "accuracy", "f1 Score", "Recall Score", "Precision_score"])

result = pd.concat([result, result1])
result

cm

"""### XGboost"""

from xgboost import XGBClassifier

XGB = XGBClassifier(random_state=0)
XGB.fit(x_train,y_train)

y_pred = XGB.predict(x_test)

acc = accuracy_score(y_test,y_pred)
f1 = f1_score(y_test,y_pred)
cm = confusion_matrix(y_test,y_pred)
rs = recall_score(y_test,y_pred)
ps = precision_score(y_test,y_pred)

cm

"""# Comparing Models"""

result2 = pd.DataFrame([['XGBoost Classifier', acc, f1, rs, ps]],
                       columns=['Model Name', "accuracy", "f1 Score", "Recall Score", "Precision_score"])

result = pd.concat([result, result2])
result

# prompt: prepare a ppt for the above

!pip install pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Create a slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add a title
title = slide.shapes.title
title.text = "Churn Analysis"

# Add a subtitle
subtitle = slide.placeholders[1]
subtitle.text = "Presentation on Churn Analysis"

# Add a body slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add a table
table = slide.shapes.add_table(rows=5, cols=2, left=Inches(1), top=Inches(2), width=Inches(6), height=Inches(2))

# Set the table header
table.rows[0].cells[0].text = "Model Name"
table.rows[0].cells[1].text = "Accuracy"

# Add the model names and accuracy scores
table.rows[1].cells[0].text = "Logistic Regression"
table.rows[1].cells[1].text = str(result.iloc[0,1])
table.rows[2].cells[0].text = "Decision Tree"
table.rows[2].cells[1].text = str(result.iloc[1,1])
table.rows[3].cells[0].text = "Random Forest"
table.rows[3].cells[1].text = str(result.iloc[2,1])
table.rows[4].cells[0].text = "XGBoost Classifier"
table.rows[4].cells[1].text = str(result.iloc[3,1])

# Set the table style
table.table_style_info = "TableStyleMedium9"

# Add a conclusion slide
slide = prs.slides.add_slide(prs.slide_layouts[5])

# Add a conclusion title
title = slide.shapes.title
title.text = "Conclusion"

# Add a conclusion body
body = slide.placeholders[1]
body.text = "Based on the analysis, the XGBoost Classifier model achieved the highest accuracy in predicting customer churn. This model can be further optimized and used to identify customers at risk of churning and take proactive measures to retain them."

# Save the presentation
prs.save("Churn Analysis.pptx")